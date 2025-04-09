import os
import zipfile
import tempfile
import re
import logging
import nbformat
from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader
# Try to import moviepy with error handling
try:
    import moviepy.editor
    from moviepy.editor import VideoFileClip
    HAS_MOVIEPY = True
except ImportError:
    HAS_MOVIEPY = False

# Set up logging
logger = logging.getLogger(__name__)

def process_file(file, file_type=None):
    """
    Process a file based on its type and extract text content.
    
    Args:
        file: The file object (from streamlit's file_uploader)
        file_type: Optional file type to override detection
        
    Returns:
        dict: Dictionary with success status and extracted text or error message
    """
    try:
        if file is None:
            return {"success": False, "error": "No file provided"}
        
        # Get file name and extension
        file_name = file.name
        extension = os.path.splitext(file_name)[1].lower() if not file_type else file_type
        
        logger.info(f"Processing file: {file_name} with extension: {extension}")
        
        # Handle different file types
        if extension == '.zip':
            return process_zip_file(file)
        elif extension == '.pptx':
            return process_pptx_file(file)
        elif extension == '.docx':
            return process_docx_file(file)
        elif extension == '.pdf':
            return process_pdf_file(file)
        elif extension in ['.py', '.txt']:
            return process_text_file(file)
        elif extension == '.ipynb':
            return process_jupyter_notebook(file)
        elif extension in ['.mp4', '.mov', '.avi', '.mkv']:
            return process_video_file(file)
        else:
            return {"success": False, "error": f"Unsupported file type: {extension}"}
    
    except Exception as e:
        logger.exception(f"Error processing file: {str(e)}")
        return {"success": False, "error": f"Error processing file: {str(e)}"}

def process_zip_file(zip_file):
    """
    Process a ZIP file containing multiple study materials.
    
    Args:
        zip_file: The ZIP file object
        
    Returns:
        dict: Dictionary with success status and extracted content or error message
    """
    try:
        # Create a temporary directory
        with tempfile.TemporaryDirectory() as temp_dir:
            # Write the zip file to the temporary directory
            zip_path = os.path.join(temp_dir, 'archive.zip')
            with open(zip_path, 'wb') as f:
                f.write(zip_file.getvalue())
            
            # Extract the contents
            combined_text = ""
            file_count = 0
            
            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                # Get list of files
                files = zip_ref.namelist()
                
                # Process each file
                for file_name in files:
                    # Skip directories and hidden files
                    if file_name.endswith('/') or file_name.startswith('.'):
                        continue
                    
                    # Get file extension
                    extension = os.path.splitext(file_name)[1].lower()
                    
                    # Extract file to temp directory
                    zip_ref.extract(file_name, temp_dir)
                    extracted_path = os.path.join(temp_dir, file_name)
                    
                    # Process based on file type
                    if extension in ['.txt', '.py']:
                        with open(extracted_path, 'r', encoding='utf-8', errors='ignore') as f:
                            file_content = f.read()
                            combined_text += f"\n\n# FILE: {file_name}\n{file_content}"
                            file_count += 1
                    
                    elif extension == '.docx':
                        try:
                            doc = Document(extracted_path)
                            file_content = '\n'.join([para.text for para in doc.paragraphs])
                            combined_text += f"\n\n# DOCUMENT: {file_name}\n{file_content}"
                            file_count += 1
                        except Exception as e:
                            logger.warning(f"Error processing DOCX file {file_name}: {str(e)}")
                    
                    elif extension == '.pptx':
                        try:
                            prs = Presentation(extracted_path)
                            file_content = '\n'.join([shape.text for slide in prs.slides 
                                                for shape in slide.shapes if hasattr(shape, "text")])
                            combined_text += f"\n\n# PRESENTATION: {file_name}\n{file_content}"
                            file_count += 1
                        except Exception as e:
                            logger.warning(f"Error processing PPTX file {file_name}: {str(e)}")
                    
                    elif extension == '.pdf':
                        try:
                            reader = PdfReader(extracted_path)
                            file_content = ""
                            for page in reader.pages:
                                file_content += page.extract_text() + "\n"
                            combined_text += f"\n\n# PDF: {file_name}\n{file_content}"
                            file_count += 1
                        except Exception as e:
                            logger.warning(f"Error processing PDF file {file_name}: {str(e)}")
                    
                    elif extension == '.ipynb':
                        try:
                            with open(extracted_path, 'r', encoding='utf-8') as f:
                                nb = nbformat.read(f, as_version=4)
                            
                            file_content = ""
                            for cell in nb.cells:
                                if cell.cell_type == 'markdown':
                                    file_content += f"# Markdown\n{cell.source}\n\n"
                                elif cell.cell_type == 'code':
                                    file_content += f"# Code\n{cell.source}\n\n"
                            
                            combined_text += f"\n\n# JUPYTER NOTEBOOK: {file_name}\n{file_content}"
                            file_count += 1
                        except Exception as e:
                            logger.warning(f"Error processing Jupyter notebook {file_name}: {str(e)}")
            
            if file_count == 0:
                return {"success": False, "error": "No valid files found in the ZIP archive"}
            
            return {"success": True, "text": combined_text, "file_count": file_count}
    
    except Exception as e:
        logger.exception(f"Error processing ZIP file: {str(e)}")
        return {"success": False, "error": f"Error processing ZIP file: {str(e)}"}

def process_pptx_file(pptx_file):
    """
    Extract text from a PowerPoint presentation.
    
    Args:
        pptx_file: The PPTX file object
        
    Returns:
        dict: Dictionary with success status and extracted text or error message
    """
    try:
        # Create a temporary file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            temp_file.write(pptx_file.getvalue())
            temp_path = temp_file.name
        
        # Extract text from the presentation
        prs = Presentation(temp_path)
        
        # Get text from slides
        text_content = ""
        for i, slide in enumerate(prs.slides):
            text_content += f"Slide {i+1}:\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content += f"{shape.text}\n"
            
            text_content += "\n"
        
        # Delete the temporary file
        os.unlink(temp_path)
        
        if not text_content.strip():
            return {"success": False, "error": "No text content found in the presentation"}
            
        return {"success": True, "text": text_content}
    
    except Exception as e:
        logger.exception(f"Error processing PowerPoint file: {str(e)}")
        return {"success": False, "error": f"Error processing PowerPoint file: {str(e)}"}

def process_docx_file(docx_file):
    """
    Extract text from a Word document.
    
    Args:
        docx_file: The DOCX file object
        
    Returns:
        dict: Dictionary with success status and extracted text or error message
    """
    try:
        # Create a temporary file
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as temp_file:
            temp_file.write(docx_file.getvalue())
            temp_path = temp_file.name
        
        # Extract text from the document
        doc = Document(temp_path)
        
        # Get text from paragraphs
        paragraphs = [para.text for para in doc.paragraphs]
        text_content = '\n'.join(paragraphs)
        
        # Get text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text for cell in row.cells)
                text_content += f"\n{row_text}"
        
        # Delete the temporary file
        os.unlink(temp_path)
        
        if not text_content.strip():
            return {"success": False, "error": "No text content found in the document"}
            
        return {"success": True, "text": text_content}
    
    except Exception as e:
        logger.exception(f"Error processing Word file: {str(e)}")
        return {"success": False, "error": f"Error processing Word file: {str(e)}"}

def process_pdf_file(pdf_file):
    """
    Extract text from a PDF document.
    
    Args:
        pdf_file: The PDF file object
        
    Returns:
        dict: Dictionary with success status and extracted text or error message
    """
    try:
        # Create a temporary file
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_file:
            temp_file.write(pdf_file.getvalue())
            temp_path = temp_file.name
        
        # Extract text from the PDF
        reader = PdfReader(temp_path)
        
        # Get text from pages
        text_content = ""
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_content += f"Page {i+1}:\n{page_text}\n\n"
        
        # Delete the temporary file
        os.unlink(temp_path)
        
        if not text_content.strip():
            return {"success": False, "error": "No text content found in the PDF"}
            
        return {"success": True, "text": text_content}
    
    except Exception as e:
        logger.exception(f"Error processing PDF file: {str(e)}")
        return {"success": False, "error": f"Error processing PDF file: {str(e)}"}

def process_text_file(text_file):
    """
    Process a plain text or Python file.
    
    Args:
        text_file: The text file object
        
    Returns:
        dict: Dictionary with success status and extracted text or error message
    """
    try:
        # Read the text file
        text_content = text_file.getvalue().decode('utf-8')
        
        if not text_content.strip():
            return {"success": False, "error": "Empty text file"}
            
        return {"success": True, "text": text_content}
    
    except Exception as e:
        logger.exception(f"Error processing text file: {str(e)}")
        return {"success": False, "error": f"Error processing text file: {str(e)}"}

def process_jupyter_notebook(ipynb_file):
    """
    Extract code and markdown content from a Jupyter notebook.
    
    Args:
        ipynb_file: The iPython notebook file object
        
    Returns:
        dict: Dictionary with success status and extracted text or error message
    """
    try:
        # Read the notebook content
        notebook_content = ipynb_file.getvalue().decode('utf-8')
        
        # Parse the notebook
        nb = nbformat.reads(notebook_content, as_version=4)
        
        # Extract text from cells
        text_content = ""
        for i, cell in enumerate(nb.cells):
            if cell.cell_type == 'markdown':
                text_content += f"Cell {i+1} (Markdown):\n{cell.source}\n\n"
            elif cell.cell_type == 'code':
                text_content += f"Cell {i+1} (Code):\n{cell.source}\n\n"
                # Include outputs if available
                if hasattr(cell, 'outputs') and cell.outputs:
                    for output in cell.outputs:
                        if output.output_type == 'stream' and 'text' in output:
                            text_content += f"Output:\n{output.text}\n"
                        elif output.output_type == 'execute_result' and 'data' in output:
                            if 'text/plain' in output.data:
                                text_content += f"Result:\n{output.data['text/plain']}\n"
        
        if not text_content.strip():
            return {"success": False, "error": "No content found in the Jupyter notebook"}
            
        return {"success": True, "text": text_content}
    
    except Exception as e:
        logger.exception(f"Error processing Jupyter notebook: {str(e)}")
        return {"success": False, "error": f"Error processing Jupyter notebook: {str(e)}"}

def process_video_file(video_file):
    """
    Extract audio from a video file and prepare it for transcription.
    
    Args:
        video_file: The video file object
        
    Returns:
        dict: Dictionary with success status and the audio file object for transcription
    """
    # Check if moviepy is available
    if not HAS_MOVIEPY:
        return {
            "success": False, 
            "error": "Video processing is not available. The moviepy module is not installed correctly."
        }
    
    try:
        # Create a temporary file for the video
        with tempfile.NamedTemporaryFile(suffix='.mp4', delete=False) as temp_video_file:
            temp_video_file.write(video_file.getvalue())
            video_path = temp_video_file.name
        
        # Create a temporary file for the extracted audio
        audio_path = video_path + ".mp3"
        
        # Extract audio from video
        logger.info(f"Extracting audio from video to {audio_path}")
        
        # Since we checked HAS_MOVIEPY at the beginning, we know VideoFileClip is available
        video = VideoFileClip(video_path)
        
        # Check if audio track exists
        if video.audio is None:
            # Clean up temporary file
            os.unlink(video_path)
            video.close()
            return {"success": False, "error": "No audio track found in the video file"}
        
        video.audio.write_audiofile(audio_path, logger=None)
        
        # Close the video file to release resources
        video.close()
        
        # Read the audio file
        with open(audio_path, 'rb') as audio_file:
            audio_data = audio_file.read()
        
        # Clean up temporary files
        os.unlink(video_path)
        os.unlink(audio_path)
        
        # Create a file-like object for the audio data
        import io
        audio_bytes = io.BytesIO(audio_data)
        audio_bytes.name = "extracted_audio.mp3"
        
        return {
            "success": True, 
            "audio_file": audio_bytes, 
            "message": "Audio extracted from video file for transcription"
        }
    
    except Exception as e:
        logger.exception(f"Error processing video file: {str(e)}")
        return {"success": False, "error": f"Error processing video file: {str(e)}"}